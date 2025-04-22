import streamlit as st
from gtts import gTTS
import pdfplumber
import pytesseract
from PIL import Image
import io
from pptx import Presentation
from docx import Document
from langdetect import detect, LangDetectException

# Optional: Set path to Tesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Page settings
st.set_page_config(
    page_title="EduVision - Learn Without Limits",
    page_icon="📚",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Logo
logo_path = "assets/eduvision_logo.png"
st.markdown("<h1 style='text-align: center; color: #4CAF50;'>EduVision</h1>", unsafe_allow_html=True)
st.markdown("<h4 style='text-align: center; color: grey;'>Empowering Visually Impaired Students Through AI</h4>", unsafe_allow_html=True)
try:
    logo = Image.open(logo_path)
    st.image(logo, width=150)
except FileNotFoundError:
    pass
st.markdown("---")

# Upload
st.subheader("📁 Upload Educational Content")
uploaded_file = st.file_uploader("Choose a file (PDF, PPTX, DOCX, TXT, PNG, JPG, JPEG)", type=['pdf', 'ppt', 'pptx', 'docx', 'txt', 'png', 'jpg', 'jpeg'])

# Supported languages
lang_options = {
    "English": "en",
    "Hindi": "hi",
    "Spanish": "es",
    "French": "fr",
    "German": "de",
    "Tamil": "ta",
    "Telugu": "te"
}
reverse_lang_map = {v: k for k, v in lang_options.items()}

# Extract text
def extract_text(file):
    text = ""
    try:
        if file.type == "application/pdf":
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
        elif file.type == "text/plain":
            text = file.read().decode("utf-8")
        elif file.type in ["image/png", "image/jpeg", "image/jpg"]:
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        st.error(f"❌ Error processing file: {e}")
    return text.strip()

# Audio generation with adjustable speed and pitch
def generate_audio(text, lang='en', speed=1.0, pitch=1.0):
    tts = gTTS(text=text, lang=lang)
    audio_io = io.BytesIO()
    tts.write_to_fp(audio_io)
    audio_io.seek(0)
    return audio_io

# Session state for recent uploads and user preferences
if "recent_uploads" not in st.session_state:
    st.session_state.recent_uploads = []

# Store user language preference in session
if "preferred_language" not in st.session_state:
    st.session_state.preferred_language = "en"

# Allow user to change language and save preference
manual_lang = st.selectbox("Choose Language", list(lang_options.keys()), index=list(lang_options.keys()).index(reverse_lang_map.get(st.session_state.preferred_language, "English")))
st.session_state.preferred_language = lang_options[manual_lang]

# Display recent uploads
st.subheader("📎 Recent Uploads")
if uploaded_file and uploaded_file.name not in st.session_state.recent_uploads:
    # Only add the file if it is not already in the recent uploads list
    st.session_state.recent_uploads.insert(0, uploaded_file.name)

for file in st.session_state.recent_uploads[:5]:  # Show latest 5 uploads
    st.write(file)

if uploaded_file:
    with st.spinner("🔍 Extracting text..."):
        extracted_text = extract_text(uploaded_file)

    if extracted_text:
        st.success("✅ Text extracted successfully!")
        st.text_area("📄 Extracted Text", extracted_text, height=200)

        # Word count and language stats
        word_count = len(extracted_text.split())
        st.markdown(f"📊 **Word Count:** {word_count} words")

        # Language detection
        try:
            detected_lang_code = detect(extracted_text)
            detected_lang_name = reverse_lang_map.get(detected_lang_code, "Unknown")
        except LangDetectException:
            detected_lang_code = "en"
            detected_lang_name = "English"

        st.markdown(f"🌐 **Detected Language:** `{detected_lang_name}`")

        # Allow user to override detected language
        override = st.checkbox("🔧 Manually choose a different language")
        if override:
            manual_lang = st.selectbox("Choose language", list(lang_options.keys()))
            selected_lang = lang_options[manual_lang]
        else:
            selected_lang = detected_lang_code

        # User controls for audio speed and pitch
        speed = st.slider("🔄 Audio Speed", min_value=0.5, max_value=2.0, value=1.0, step=0.1)
        pitch = st.slider("🎵 Audio Pitch", min_value=0.5, max_value=2.0, value=1.0, step=0.1)

        # Audio Generation Button
        if st.button("🔊 Convert to Audio"):
            with st.spinner(f"🎙️ Generating audio in {manual_lang}..."):
                audio_file = generate_audio(extracted_text, lang=selected_lang, speed=speed, pitch=pitch)
                st.audio(audio_file, format="audio/mp3")
                st.download_button("⬇️ Download Audio", data=audio_file, file_name="eduvision_audio.mp3")
            st.success("🎧 Audio generated successfully!")

    else:
        st.error("❌ Couldn't extract any text. Please try another file.")

# Footer
st.markdown("---")
st.markdown("<p style='text-align: center; color: grey;'>Made with ❤️ for AlgoArena</p>", unsafe_allow_html=True)
